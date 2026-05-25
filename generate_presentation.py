"""
ORCHIIMMO - Générateur de Présentation PPTX pour la Soutenance SFE
Style: Dark theme, couleurs #4895EF + #7C3AED + #06D6A0
Transitions: Morph + Fade + Wipe
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Cm
from pptx.enum.dml import MSO_THEME_COLOR
from lxml import etree
import copy

# ─────────────────────────────────────────────────────────────────────────────
# PALETTE ORCHIIMMO
# ─────────────────────────────────────────────────────────────────────────────
BG_DARK       = RGBColor(0x08, 0x0B, 0x14)   # #080B14 fond principal
BG_CARD       = RGBColor(0x12, 0x17, 0x28)   # #121728 cartes
BG_SURFACE    = RGBColor(0x1A, 0x1F, 0x35)   # #1A1F35 surface
BLUE          = RGBColor(0x48, 0x95, 0xEF)   # #4895EF bleu principal
PURPLE        = RGBColor(0x7C, 0x3A, 0xED)   # #7C3AED violet
TEAL          = RGBColor(0x06, 0xD6, 0xA0)   # #06D6A0 vert/teal
ORANGE        = RGBColor(0xF4, 0xA2, 0x61)   # #F4A261 orange
RED           = RGBColor(0xEF, 0x44, 0x44)   # #EF4444 rouge
YELLOW        = RGBColor(0xF5, 0x9E, 0x0B)   # #F59E0B jaune
TEXT_LIGHT    = RGBColor(0xE2, 0xE8, 0xF0)   # #E2E8F0 texte clair
TEXT_MUTED    = RGBColor(0x88, 0x92, 0xA4)   # #8892A4 texte secondaire
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
BLACK         = RGBColor(0x00, 0x00, 0x00)

# ─────────────────────────────────────────────────────────────────────────────
# DIMENSIONS (16:9 widescreen)
# ─────────────────────────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK_LAYOUT = prs.slide_layouts[6]   # blank

# ─────────────────────────────────────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color=None, alpha=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, font_size=18, bold=False, italic=False,
             color=TEXT_LIGHT, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_multiline_text(slide, lines, x, y, w, h, font_size=18, bold=False,
                        color=TEXT_LIGHT, align=PP_ALIGN.LEFT, spacing=None):
    """lines = list of (text, color, size, bold) or just strings"""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            t, c, s, b = item, color, font_size, bold
        else:
            t = item.get('text', '')
            c = item.get('color', color)
            s = item.get('size', font_size)
            b = item.get('bold', bold)

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.space_before = Pt(spacing)
        run = p.add_run()
        run.text = t
        run.font.size = Pt(s)
        run.font.bold = b
        run.font.color.rgb = c
        run.font.name = "Segoe UI"
    return txBox


def add_gradient_rect(slide, x, y, w, h, color1, color2, angle=0):
    """Gradient via two overlapping semitransparent shapes."""
    shape = slide.shapes.add_shape(1, x, y, w, h)
    # Use gradient via XML
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color1
    shape.line.fill.background()

    sp = shape._element
    spPr = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')
    if spPr is None:
        return shape

    # Build gradient fill XML
    nsmap = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    gradFill = etree.SubElement(spPr, f'{{{nsmap}}}gradFill')
    gsLst = etree.SubElement(gradFill, f'{{{nsmap}}}gsLst')

    gs1 = etree.SubElement(gsLst, f'{{{nsmap}}}gs')
    gs1.set('pos', '0')
    srgb1 = etree.SubElement(gs1, f'{{{nsmap}}}srgbClr')
    srgb1.set('val', f'{color1.red:02X}{color1.green:02X}{color1.blue:02X}')

    gs2 = etree.SubElement(gsLst, f'{{{nsmap}}}gs')
    gs2.set('pos', '100000')
    srgb2 = etree.SubElement(gs2, f'{{{nsmap}}}srgbClr')
    srgb2.set('val', f'{color2.red:02X}{color2.green:02X}{color2.blue:02X}')

    lin = etree.SubElement(gradFill, f'{{{nsmap}}}lin')
    lin.set('ang', str(angle * 60000))
    lin.set('scaled', '0')

    # Replace solidFill with gradFill
    solidFill = spPr.find(f'{{{nsmap}}}solidFill')
    if solidFill is not None:
        spPr.remove(solidFill)

    return shape


def set_bg(slide, color=BG_DARK):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_morph_transition(slide):
    """Add Morph transition to a slide."""
    transition_xml = (
        '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" '
        'spd="slow" advTm="0">'
        '<p14:transition dur="800" invX="0" invY="0">'
        '<p14:morph/>'
        '</p14:transition>'
        '</p:transition>'
    )
    elem = etree.fromstring(transition_xml)
    slide._element.append(elem)


def add_fade_transition(slide, dur=600):
    transition_xml = (
        f'<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        f'spd="slow" advTm="0">'
        f'<p:fade/>'
        f'</p:transition>'
    )
    elem = etree.fromstring(transition_xml)
    slide._element.append(elem)


def add_push_transition(slide, dir="l"):
    """Push transition: l=left, r=right, u=up, d=down."""
    transition_xml = (
        f'<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        f'spd="slow" advTm="0">'
        f'<p:push dir="{dir}"/>'
        f'</p:transition>'
    )
    elem = etree.fromstring(transition_xml)
    slide._element.append(elem)


def add_wipe_transition(slide, dir="l"):
    transition_xml = (
        f'<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        f'spd="slow" advTm="0">'
        f'<p:wipe dir="{dir}"/>'
        f'</p:transition>'
    )
    elem = etree.fromstring(transition_xml)
    slide._element.append(elem)


def add_reveal_transition(slide, dir="l"):
    transition_xml = (
        f'<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        f'spd="slow" advTm="0">'
        f'<p:cover dir="{dir}"/>'
        f'</p:transition>'
    )
    elem = etree.fromstring(transition_xml)
    slide._element.append(elem)


def section_header_bar(slide, title, accent_color=BLUE):
    """Accent top bar + section title."""
    add_rect(slide, 0, 0, W, Inches(0.08), fill_color=accent_color)
    # Left accent stripe
    add_rect(slide, Inches(0.5), Inches(0.15), Inches(0.06), Inches(0.5), fill_color=accent_color)
    add_text(slide, title, Inches(0.65), Inches(0.12), Inches(11), Inches(0.65),
             font_size=28, bold=True, color=TEXT_LIGHT)


def add_card(slide, x, y, w, h, title=None, content_lines=None,
             accent=BLUE, icon=None):
    """Glassmorphism-style card."""
    card = add_rect(slide, x, y, w, h, fill_color=BG_CARD)
    # Left accent border
    add_rect(slide, x, y, Inches(0.045), h, fill_color=accent)
    if icon and title:
        add_text(slide, f"{icon}  {title}", x + Inches(0.12), y + Inches(0.12),
                 w - Inches(0.2), Inches(0.4), font_size=14, bold=True, color=accent)
        start_y = y + Inches(0.55)
    elif title:
        add_text(slide, title, x + Inches(0.12), y + Inches(0.12),
                 w - Inches(0.2), Inches(0.4), font_size=13, bold=True, color=accent)
        start_y = y + Inches(0.55)
    else:
        start_y = y + Inches(0.15)

    if content_lines:
        add_multiline_text(slide, content_lines, x + Inches(0.15), start_y,
                           w - Inches(0.3), h - (start_y - y) - Inches(0.15),
                           font_size=11, color=TEXT_LIGHT, spacing=3)
    return card


def add_progress_bar(slide, x, y, w, h, percent, bg_color=BG_SURFACE, fill_color=BLUE):
    """Horizontal progress bar."""
    add_rect(slide, x, y, w, h, fill_color=bg_color)
    add_rect(slide, x, y, int(w * percent / 100), h, fill_color=fill_color)


def add_divider(slide, y, color=BLUE, opacity_hex="40"):
    line_shape = slide.shapes.add_shape(1, Inches(0.5), y, W - Inches(1), Inches(0.012))
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = color
    line_shape.line.fill.background()


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — COUVERTURE
# ─────────────────────────────────────────────────────────────────────────────
def slide_cover():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_bg(slide, BG_DARK)

    # Background gradient overlay (large rectangles)
    add_gradient_rect(slide, 0, 0, W, H, BG_DARK, BG_CARD, angle=135)

    # Top decorative gradient bar
    add_gradient_rect(slide, 0, 0, W, Inches(0.12), BLUE, PURPLE, angle=0)

    # Bottom bar
    add_gradient_rect(slide, 0, H - Inches(0.12), W, Inches(0.12), PURPLE, BLUE, angle=0)

    # Large decorative circle (top-right)
    circ = slide.shapes.add_shape(9, Inches(9.5), Inches(-1.5), Inches(5), Inches(5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = BLUE
    sp = circ._element
    spPr = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')
    # Reduce opacity
    circ.line.fill.background()

    # Geometric accent shapes
    add_gradient_rect(slide, Inches(10.2), Inches(1.5), Inches(2.8), Inches(0.8), BLUE, PURPLE, angle=0)
    add_gradient_rect(slide, Inches(10.6), Inches(2.5), Inches(2.1), Inches(0.08), TEAL, BLUE, angle=0)

    # Small decorative dots pattern
    for i in range(5):
        dot = slide.shapes.add_shape(9, Inches(10.5 + i*0.35), Inches(4.5), Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = TEAL if i % 2 == 0 else PURPLE
        dot.line.fill.background()

    # LOGO / Brand block
    add_rect(slide, Inches(0.8), Inches(0.9), Inches(0.08), Inches(1.2), fill_color=BLUE)

    # ORCHI text (large)
    add_text(slide, "ORCHI", Inches(1.1), Inches(0.9), Inches(6), Inches(1.0),
             font_size=72, bold=True, color=BLUE, font_name="Segoe UI")

    # IMMO text (large, purple)
    add_text(slide, "IMMO", Inches(4.55), Inches(0.9), Inches(5), Inches(1.0),
             font_size=72, bold=True, color=PURPLE, font_name="Segoe UI")

    # Underline accent
    add_gradient_rect(slide, Inches(1.1), Inches(1.85), Inches(7.5), Inches(0.06), BLUE, PURPLE, angle=0)

    # Tagline
    add_text(slide, "Plateforme d'Analyse Immobilière Intelligente — 100% Maroc",
             Inches(1.1), Inches(2.0), Inches(10), Inches(0.5),
             font_size=18, bold=False, color=TEXT_MUTED, font_name="Segoe UI")

    # Divider
    add_divider(slide, Inches(2.65), BLUE)

    # Title block
    add_text(slide, "Soutenance de Stage de Fin d'Études",
             Inches(1.1), Inches(2.8), Inches(10), Inches(0.55),
             font_size=22, bold=True, color=TEXT_LIGHT)

    # Description
    add_text(slide,
             "Développement d'une plateforme web d'estimation et d'analyse\n"
             "des prix immobiliers au Maroc par Intelligence Artificielle",
             Inches(1.1), Inches(3.4), Inches(10), Inches(1.0),
             font_size=15, color=TEXT_MUTED)

    # Info cards row
    info_items = [
        ("📅", "2025 — 2026", "Année"),
        ("🎓", "PFE / SFE", "Type"),
        ("🇲🇦", "Maroc", "Périmètre"),
    ]
    for i, (icon, val, label) in enumerate(info_items):
        cx = Inches(1.1 + i * 3.0)
        add_rect(slide, cx, Inches(4.7), Inches(2.6), Inches(1.0), fill_color=BG_CARD)
        add_rect(slide, cx, Inches(4.7), Inches(0.04), Inches(1.0),
                 fill_color=[BLUE, PURPLE, TEAL][i])
        add_text(slide, icon + "  " + val, cx + Inches(0.15), Inches(4.8),
                 Inches(2.3), Inches(0.4), font_size=16, bold=True,
                 color=[BLUE, PURPLE, TEAL][i])
        add_text(slide, label, cx + Inches(0.15), Inches(5.15), Inches(2.3), Inches(0.3),
                 font_size=11, color=TEXT_MUTED)

    # Bottom attribution
    add_text(slide, "hadifatimaezzahra5@gmail.com",
             Inches(0.5), Inches(6.8), Inches(8), Inches(0.4),
             font_size=11, color=TEXT_MUTED, align=PP_ALIGN.LEFT)
    add_text(slide, "Mai 2026", Inches(10.5), Inches(6.8), Inches(2.5), Inches(0.4),
             font_size=11, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)

    add_fade_transition(slide)
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — PLAN / SOMMAIRE
# ─────────────────────────────────────────────────────────────────────────────
def slide_plan():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_bg(slide, BG_DARK)
    add_gradient_rect(slide, 0, 0, W, Inches(0.08), BLUE, PURPLE, angle=0)

    # Title
    add_text(slide, "Plan de la Présentation", Inches(0.5), Inches(0.2),
             Inches(10), Inches(0.7), font_size=32, bold=True, color=TEXT_LIGHT)
    add_gradient_rect(slide, Inches(0.5), Inches(0.85), Inches(3), Inches(0.05),
                      BLUE, PURPLE, angle=0)

    items = [
        ("01", "Contexte & Problématique",        BLUE),
        ("02", "Présentation du Projet ORCHIIMMO", PURPLE),
        ("03", "Architecture Technique",           TEAL),
        ("04", "Scraping & Données",               ORANGE),
        ("05", "Catalogue & Carte Interactive",    BLUE),
        ("06", "Estimation IA — Moteur ML",        PURPLE),
        ("07", "Dashboards Power BI",              TEAL),
        ("08", "Chatbot Intelligent",              ORANGE),
        ("09", "Déploiement & Production",         RED),
        ("10", "Résultats & Perspectives",         YELLOW),
    ]

    cols = [items[:5], items[5:]]
    for col_idx, col_items in enumerate(cols):
        cx = Inches(0.5 + col_idx * 6.4)
        for i, (num, label, color) in enumerate(col_items):
            cy = Inches(1.2 + i * 1.1)
            # card bg
            add_rect(slide, cx, cy, Inches(5.9), Inches(0.9), fill_color=BG_CARD)
            # number badge
            add_rect(slide, cx, cy, Inches(0.55), Inches(0.9), fill_color=color)
            add_text(slide, num, cx + Inches(0.02), cy + Inches(0.18),
                     Inches(0.5), Inches(0.5), font_size=18, bold=True,
                     color=WHITE, align=PP_ALIGN.CENTER)
            # label
            add_text(slide, label, cx + Inches(0.7), cy + Inches(0.22),
                     Inches(5.0), Inches(0.45), font_size=14, bold=False, color=TEXT_LIGHT)

    add_morph_transition(slide)
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — CONTEXTE & PROBLEMATIQUE
# ─────────────────────────────────────────────────────────────────────────────
def slide_contexte():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_bg(slide, BG_DARK)
    section_header_bar(slide, "01 — Contexte & Problématique", BLUE)

    # Left column - problem
    add_rect(slide, Inches(0.4), Inches(1.1), Inches(6.0), Inches(5.8), fill_color=BG_CARD)
    add_rect(slide, Inches(0.4), Inches(1.1), Inches(0.06), Inches(5.8), fill_color=RED)

    add_text(slide, "⚠  Problématique", Inches(0.6), Inches(1.2),
             Inches(5.6), Inches(0.45), font_size=16, bold=True, color=RED)

    problems = [
        {"text": "Opacité des prix", "color": TEXT_LIGHT, "size": 13, "bold": True},
        {"text": "Aucune donnée agrégée et fiable sur les prix du marché marocain", "color": TEXT_MUTED, "size": 11, "bold": False},
        {"text": "", "color": TEXT_MUTED, "size": 6, "bold": False},
        {"text": "Fragmentation des sources", "color": TEXT_LIGHT, "size": 13, "bold": True},
        {"text": "Données dispersées sur 9+ portails immobiliers différents", "color": TEXT_MUTED, "size": 11, "bold": False},
        {"text": "", "color": TEXT_MUTED, "size": 6, "bold": False},
        {"text": "Pas d'outil d'estimation", "color": TEXT_LIGHT, "size": 13, "bold": True},
        {"text": "Absence d'un moteur IA local pour prédire les prix au Maroc", "color": TEXT_MUTED, "size": 11, "bold": False},
        {"text": "", "color": TEXT_MUTED, "size": 6, "bold": False},
        {"text": "Décision difficile", "color": TEXT_LIGHT, "size": 13, "bold": True},
        {"text": "Acheteurs et investisseurs manquent de visibilité analytique", "color": TEXT_MUTED, "size": 11, "bold": False},
    ]
    add_multiline_text(slide, problems, Inches(0.65), Inches(1.75),
                       Inches(5.6), Inches(4.8), spacing=2)

    # Right column - solution / stats
    add_rect(slide, Inches(6.8), Inches(1.1), Inches(6.0), Inches(5.8), fill_color=BG_CARD)
    add_rect(slide, Inches(6.8), Inches(1.1), Inches(0.06), Inches(5.8), fill_color=TEAL)

    add_text(slide, "📊  Marché Immobilier Marocain", Inches(7.0), Inches(1.2),
             Inches(5.6), Inches(0.45), font_size=16, bold=True, color=TEAL)

    # KPI stats
    kpi_items = [
        ("7 805", "Annonces actives collectées", BLUE),
        ("331", "Villes couvertes", PURPLE),
        ("1 350 000 MAD", "Prix médian national", TEAL),
        ("9 sources", "Portails immobiliers scrappés", ORANGE),
    ]
    for i, (val, label, color) in enumerate(kpi_items):
        ky = Inches(1.85 + i * 1.2)
        add_rect(slide, Inches(7.0), ky, Inches(5.5), Inches(1.0), fill_color=BG_SURFACE)
        add_rect(slide, Inches(7.0), ky, Inches(0.04), Inches(1.0), fill_color=color)
        add_text(slide, val, Inches(7.15), ky + Inches(0.1),
                 Inches(4.0), Inches(0.45), font_size=22, bold=True, color=color)
        add_text(slide, label, Inches(7.15), ky + Inches(0.52),
                 Inches(4.8), Inches(0.3), font_size=11, color=TEXT_MUTED)

    add_morph_transition(slide)
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — PRESENTATION DU PROJET
# ─────────────────────────────────────────────────────────────────────────────
def slide_projet():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_bg(slide, BG_DARK)
    section_header_bar(slide, "02 — Présentation du Projet ORCHIIMMO", PURPLE)

    # Big logo area
    add_gradient_rect(slide, Inches(0.4), Inches(1.0), Inches(12.5), Inches(1.2),
                      BG_CARD, BG_SURFACE, angle=0)
    add_rect(slide, Inches(0.4), Inches(1.0), Inches(0.08), Inches(1.2), fill_color=PURPLE)
    add_text(slide, "ORCHI", Inches(0.65), Inches(1.05), Inches(3.5), Inches(1.0),
             font_size=48, bold=True, color=BLUE)
    add_text(slide, "IMMO", Inches(3.0), Inches(1.05), Inches(4), Inches(1.0),
             font_size=48, bold=True, color=PURPLE)
    add_text(slide, "— Plateforme d'analyse immobilière 100% Maroc",
             Inches(6.4), Inches(1.38), Inches(6.5), Inches(0.5),
             font_size=14, color=TEXT_MUTED)

    # Feature cards
    features = [
        ("🏠", "Catalogue\nImmobilier", "Recherche avancée\n7 805+ annonces\nFiltres multicritères\nGéolocalisation", BLUE),
        ("🤖", "Estimation\npar IA", "4 modèles ML\nLightGBM · XGBoost\nCatBoost · RandomForest\nIntervalle de confiance", PURPLE),
        ("📊", "Dashboards\nPower BI", "Analyses de marché\nTendances prix\nKPIs immobiliers\nVisualisation interactive", TEAL),
        ("💬", "Chatbot\nIntelligent", "Assistant virtuel\nIntégration N8N\nFrançais / Arabe\nAvailable 24/7", ORANGE),
    ]
    for i, (icon, title, desc, color) in enumerate(features):
        cx = Inches(0.4 + i * 3.15)
        cy = Inches(2.5)
        ch = Inches(4.5)
        cw = Inches(3.0)

        add_rect(slide, cx, cy, cw, ch, fill_color=BG_CARD)
        add_gradient_rect(slide, cx, cy, cw, Inches(0.06), color, PURPLE, angle=0)

        # Icon circle
        ic = slide.shapes.add_shape(9, cx + Inches(0.9), cy + Inches(0.25),
                                     Inches(1.2), Inches(1.2))
        ic.fill.solid()
        ic.fill.fore_color.rgb = BG_SURFACE
        ic.line.color.rgb = color
        ic.line.width = Pt(1.5)

        add_text(slide, icon, cx + Inches(0.9), cy + Inches(0.35),
                 Inches(1.2), Inches(1.0), font_size=32, align=PP_ALIGN.CENTER)

        add_text(slide, title, cx + Inches(0.1), cy + Inches(1.55),
                 Inches(2.8), Inches(0.7), font_size=14, bold=True,
                 color=color, align=PP_ALIGN.CENTER)
        add_text(slide, desc, cx + Inches(0.15), cy + Inches(2.35),
                 Inches(2.7), Inches(1.7), font_size=10.5, color=TEXT_MUTED,
                 align=PP_ALIGN.CENTER)

    add_reveal_transition(slide, "l")
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — ARCHITECTURE TECHNIQUE
# ─────────────────────────────────────────────────────────────────────────────
def slide_architecture():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_bg(slide, BG_DARK)
    section_header_bar(slide, "03 — Architecture Technique", TEAL)

    # Stack layers
    layers = [
        ("FRONTEND",   "Bootstrap 5.3 · CSS Glassmorphism · Vanilla JS · Bilingual (FR/AR)", BLUE),
        ("BACKEND",    "Django 5.2 · Python 3.12 · WhiteNoise · Gunicorn WSGI",             PURPLE),
        ("IA / ML",    "LightGBM · XGBoost · CatBoost · RandomForest · scikit-learn",        TEAL),
        ("DONNÉES",    "PostgreSQL · BeautifulSoup4 · Requests · lxml · pandas/numpy",        ORANGE),
        ("SERVICES",   "N8N Chatbot · Power BI Embed · Brevo Email · Pillow · openpyxl",     YELLOW),
        ("DÉPLOIEMENT","Render.com · Frankfurt · Gunicorn · WhiteNoise · django-ratelimit",   RED),
    ]

    for i, (layer, tech, color) in enumerate(layers):
        cy = Inches(1.0 + i * 1.05)
        # Full-width card
        add_rect(slide, Inches(0.4), cy, Inches(12.5), Inches(0.88), fill_color=BG_CARD)
        add_rect(slide, Inches(0.4), cy, Inches(0.06), Inches(0.88), fill_color=color)
        # Layer label badge
        add_rect(slide, Inches(0.5), cy + Inches(0.18), Inches(1.6), Inches(0.52),
                 fill_color=color)
        add_text(slide, layer, Inches(0.5), cy + Inches(0.22),
                 Inches(1.6), Inches(0.44), font_size=11, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)
        # Tech description
        add_text(slide, tech, Inches(2.3), cy + Inches(0.22),
                 Inches(10.4), Inches(0.44), font_size=12, color=TEXT_LIGHT)

    # Arrows / connectors (decorative)
    for i in range(5):
        cy = Inches(1.87 + i * 1.05)
        add_text(slide, "▼", Inches(6.1), cy, Inches(0.5), Inches(0.18),
                 font_size=10, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    add_morph_transition(slide)
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — SCRAPING & DONNÉES
# ─────────────────────────────────────────────────────────────────────────────
def slide_scraping():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_bg(slide, BG_DARK)
    section_header_bar(slide, "04 — Collecte de Données — Web Scraping", ORANGE)

    # Sources grid
    sources = [
        ("Mubawab.ma",       "#1"),
        ("Avito.ma",         "#2"),
        ("Sarouty.ma",       "#3"),
        ("Agenz.ma",         "#4"),
        ("MarocAnnonces.ma", "#5"),
        ("Masaken.ma",       "#6"),
        ("LogicImmo.ma",     "#7"),
        ("Bikhir.ma",        "#8"),
        ("+ Sources",        "#9"),
    ]
    colors_cycle = [BLUE, PURPLE, TEAL, ORANGE, RED, YELLOW, BLUE, PURPLE, TEAL]

    add_text(slide, "9 Portails Immobiliers Scrappés", Inches(0.4), Inches(1.0),
             Inches(7), Inches(0.45), font_size=18, bold=True, color=TEXT_LIGHT)

    for i, ((name, num), color) in enumerate(zip(sources, colors_cycle)):
        row = i // 3
        col = i % 3
        cx = Inches(0.4 + col * 2.3)
        cy = Inches(1.5 + row * 1.35)
        add_rect(slide, cx, cy, Inches(2.1), Inches(1.1), fill_color=BG_CARD)
        add_rect(slide, cx, cy, Inches(0.06), Inches(1.1), fill_color=color)
        add_text(slide, num, cx + Inches(0.15), cy + Inches(0.05),
                 Inches(0.5), Inches(0.35), font_size=11, bold=True, color=color)
        add_text(slide, name, cx + Inches(0.15), cy + Inches(0.38),
                 Inches(1.85), Inches(0.5), font_size=11, bold=False, color=TEXT_LIGHT)

    # Right side — pipeline
    add_rect(slide, Inches(7.4), Inches(1.0), Inches(5.5), Inches(6.2), fill_color=BG_CARD)
    add_rect(slide, Inches(7.4), Inches(1.0), Inches(0.06), Inches(6.2), fill_color=ORANGE)
    add_text(slide, "🔄  Pipeline de Collecte", Inches(7.6), Inches(1.1),
             Inches(5.2), Inches(0.45), font_size=15, bold=True, color=ORANGE)

    pipeline_steps = [
        ("1", "Requêtes HTTP (requests)", BLUE),
        ("2", "Parsing HTML (BeautifulSoup4 + lxml)", PURPLE),
        ("3", "Extraction: titre, prix, surface, ville", TEAL),
        ("4", "Nettoyage & normalisation des données", ORANGE),
        ("5", "Géocodage (lat/lng par ville)", YELLOW),
        ("6", "Stockage PostgreSQL / SQLite", RED),
        ("7", "Scraping automatique configurable", BLUE),
    ]
    for i, (step, desc, color) in enumerate(pipeline_steps):
        sy = Inches(1.65 + i * 0.76)
        add_rect(slide, Inches(7.6), sy, Inches(4.9), Inches(0.62), fill_color=BG_SURFACE)
        # Step badge
        add_rect(slide, Inches(7.6), sy, Inches(0.4), Inches(0.62), fill_color=color)
        add_text(slide, step, Inches(7.6), sy + Inches(0.12),
                 Inches(0.4), Inches(0.35), font_size=12, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, desc, Inches(8.1), sy + Inches(0.14),
                 Inches(4.3), Inches(0.35), font_size=11, color=TEXT_LIGHT)

    add_wipe_transition(slide, "l")
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — CATALOGUE & RECHERCHE
# ─────────────────────────────────────────────────────────────────────────────
def slide_catalogue():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_bg(slide, BG_DARK)
    section_header_bar(slide, "05 — Catalogue Immobilier & Carte Interactive", BLUE)

    # Main feature blocks
    features = [
        ("🔍", "Recherche Avancée",
         "• Filtres: ville, type, prix min/max\n"
         "• Surface (m²), nbre chambres/SDB\n"
         "• Tri par prix, surface, date\n"
         "• Pagination intelligente",
         BLUE),
        ("🗺️", "Carte Interactive",
         "• Géolocalisation de chaque bien\n"
         "• Coordonnées lat/lng par ville\n"
         "• Clustering des marqueurs\n"
         "• Vue satellite disponible",
         TEAL),
        ("📱", "Fiche Propriété",
         "• Photos HD de l'annonce\n"
         "• Prix en MAD + Prix/m²\n"
         "• Contact WhatsApp direct\n"
         "• Source originale linkée",
         PURPLE),
        ("📈", "Statistiques Marché",
         "• Top 10 villes (Marrakech, Casa...)\n"
         "• Prix médian: 1 350 000 MAD\n"
         "• Prix moyen: 2 837 789 MAD\n"
         "• 331 villes couvertes",
         ORANGE),
    ]

    for i, (icon, title, desc, color) in enumerate(features):
        row = i // 2
        col = i % 2
        cx = Inches(0.4 + col * 6.45)
        cy = Inches(1.0 + row * 2.9)
        cw = Inches(6.1)
        ch = Inches(2.65)

        add_rect(slide, cx, cy, cw, ch, fill_color=BG_CARD)
        add_gradient_rect(slide, cx, cy, cw, Inches(0.06), color, PURPLE, angle=0)

        add_text(slide, f"{icon}  {title}", cx + Inches(0.2), cy + Inches(0.15),
                 Inches(5.5), Inches(0.45), font_size=16, bold=True, color=color)
        add_text(slide, desc, cx + Inches(0.2), cy + Inches(0.7),
                 Inches(5.7), Inches(1.8), font_size=12, color=TEXT_MUTED)

    add_morph_transition(slide)
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — ESTIMATION IA
# ─────────────────────────────────────────────────────────────────────────────
def slide_estimation():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_bg(slide, BG_DARK)
    section_header_bar(slide, "06 — Moteur d'Estimation par IA", PURPLE)

    # Left: input form mock
    add_rect(slide, Inches(0.4), Inches(1.0), Inches(5.5), Inches(6.2), fill_color=BG_CARD)
    add_gradient_rect(slide, Inches(0.4), Inches(1.0), Inches(5.5), Inches(0.06),
                      PURPLE, BLUE, angle=0)
    add_text(slide, "📋  Formulaire d'Estimation", Inches(0.55), Inches(1.1),
             Inches(5.2), Inches(0.45), font_size=15, bold=True, color=PURPLE)

    form_fields = [
        ("🏙️  Ville", "Marrakech, Casablanca, Tanger..."),
        ("📐  Surface (m²)", "Ex: 85 m²"),
        ("🛏️  Chambres", "1 — 5+"),
        ("🚿  Salles de bain", "1 — 3+"),
        ("🏠  Type de bien", "Appartement / Villa / Studio"),
        ("📍  Quartier", "Centre-ville, Guéliz, Bourgogne..."),
    ]
    for i, (label, placeholder) in enumerate(form_fields):
        fy = Inches(1.65 + i * 0.8)
        add_text(slide, label, Inches(0.6), fy, Inches(5.0), Inches(0.3),
                 font_size=11, bold=True, color=TEXT_MUTED)
        add_rect(slide, Inches(0.6), fy + Inches(0.3), Inches(5.0), Inches(0.38),
                 fill_color=BG_SURFACE, line_color=PURPLE, line_width=Pt(0.5))
        add_text(slide, placeholder, Inches(0.7), fy + Inches(0.32),
                 Inches(4.8), Inches(0.32), font_size=10, color=TEXT_MUTED)

    # Estimate button
    add_gradient_rect(slide, Inches(0.6), Inches(6.5), Inches(5.0), Inches(0.52),
                      PURPLE, BLUE, angle=0)
    add_text(slide, "⚡  Estimer le Prix", Inches(0.6), Inches(6.52),
             Inches(5.0), Inches(0.48), font_size=14, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)

    # Right: result display
    add_rect(slide, Inches(6.3), Inches(1.0), Inches(6.6), Inches(6.2), fill_color=BG_CARD)
    add_gradient_rect(slide, Inches(6.3), Inches(1.0), Inches(6.6), Inches(0.06),
                      TEAL, BLUE, angle=0)
    add_text(slide, "📊  Résultat de l'Estimation", Inches(6.5), Inches(1.1),
             Inches(6.2), Inches(0.45), font_size=15, bold=True, color=TEAL)

    # Big price display
    add_rect(slide, Inches(6.5), Inches(1.65), Inches(6.1), Inches(1.5), fill_color=BG_SURFACE)
    add_text(slide, "1 247 500 MAD", Inches(6.5), Inches(1.75),
             Inches(6.1), Inches(0.8), font_size=36, bold=True,
             color=TEAL, align=PP_ALIGN.CENTER)
    add_text(slide, "Prix estimé par l'IA", Inches(6.5), Inches(2.5),
             Inches(6.1), Inches(0.4), font_size=12, color=TEXT_MUTED,
             align=PP_ALIGN.CENTER)

    # Confidence interval
    add_rect(slide, Inches(6.5), Inches(3.3), Inches(6.1), Inches(0.8), fill_color=BG_SURFACE)
    add_text(slide, "Intervalle de confiance: 1 100 000 — 1 395 000 MAD",
             Inches(6.6), Inches(3.4), Inches(5.9), Inches(0.5),
             font_size=11, color=ORANGE)

    # Metrics
    metrics = [
        ("Prix / m²", "14 676 MAD/m²", BLUE),
        ("Modèle", "LightGBM v2.1", PURPLE),
        ("Confiance", "±11.9%", TEAL),
        ("Ville", "Marrakech", ORANGE),
    ]
    for i, (label, val, color) in enumerate(metrics):
        my = Inches(4.2 + i * 0.72)
        add_rect(slide, Inches(6.5), my, Inches(6.1), Inches(0.58), fill_color=BG_SURFACE)
        add_rect(slide, Inches(6.5), my, Inches(0.04), Inches(0.58), fill_color=color)
        add_text(slide, label, Inches(6.65), my + Inches(0.12),
                 Inches(3.0), Inches(0.32), font_size=11, color=TEXT_MUTED)
        add_text(slide, val, Inches(9.6), my + Inches(0.12),
                 Inches(2.8), Inches(0.32), font_size=12, bold=True, color=color,
                 align=PP_ALIGN.RIGHT)

    add_morph_transition(slide)
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — MODELES ML
# ─────────────────────────────────────────────────────────────────────────────
def slide_ml():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_bg(slide, BG_DARK)
    section_header_bar(slide, "06b — Modèles Machine Learning Comparés", TEAL)

    models = [
        ("LightGBM", "🏆 Meilleur modèle", "R² Score: 0.912\nRMSE: 182 400 MAD\nTrès rapide, gestion catégorielle native\nIdéal pour données immobilières", TEAL, True),
        ("XGBoost", "2ème modèle", "R² Score: 0.891\nRMSE: 198 100 MAD\nRobuste, gradient boosting\nBonne généralisation", BLUE, False),
        ("CatBoost", "3ème modèle", "R² Score: 0.878\nRMSE: 210 500 MAD\nGestion auto des catégories\nMoins de tuning requis", PURPLE, False),
        ("RandomForest", "Baseline", "R² Score: 0.843\nRMSE: 235 800 MAD\nSimple et interprétable\nRéférence pour comparaison", ORANGE, False),
    ]

    for i, (name, badge, desc, color, is_best) in enumerate(models):
        cx = Inches(0.4 + i * 3.15)
        cy = Inches(1.1)

        # Card height larger for winner
        ch = Inches(5.7) if is_best else Inches(5.3)
        add_rect(slide, cx, cy, Inches(3.0), ch, fill_color=BG_CARD)
        add_gradient_rect(slide, cx, cy, Inches(3.0), Inches(0.06),
                          color, PURPLE, angle=0)

        if is_best:
            # Crown badge
            add_rect(slide, cx, cy - Inches(0.25), Inches(3.0), Inches(0.25),
                     fill_color=TEAL)
            add_text(slide, "🏆 MEILLEUR MODÈLE SÉLECTIONNÉ",
                     cx, cy - Inches(0.24), Inches(3.0), Inches(0.22),
                     font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        add_text(slide, name, cx + Inches(0.1), cy + Inches(0.15),
                 Inches(2.8), Inches(0.55), font_size=20, bold=True, color=color)
        add_text(slide, badge, cx + Inches(0.1), cy + Inches(0.7),
                 Inches(2.8), Inches(0.3), font_size=11, color=TEXT_MUTED)

        # Separator
        add_rect(slide, cx + Inches(0.1), cy + Inches(1.05), Inches(2.7), Inches(0.02),
                 fill_color=color)

        # R² bar
        add_text(slide, "R² Score", cx + Inches(0.1), cy + Inches(1.15),
                 Inches(2.8), Inches(0.3), font_size=10, color=TEXT_MUTED)
        r2_scores = {"LightGBM": 91, "XGBoost": 89, "CatBoost": 88, "RandomForest": 84}
        r2 = r2_scores.get(name, 80)
        add_rect(slide, cx + Inches(0.1), cy + Inches(1.45), Inches(2.7), Inches(0.22),
                 fill_color=BG_SURFACE)
        add_rect(slide, cx + Inches(0.1), cy + Inches(1.45),
                 int(Inches(2.7) * r2 / 100), Inches(0.22), fill_color=color)
        add_text(slide, f"{r2}%", cx + Inches(2.35), cy + Inches(1.45),
                 Inches(0.45), Inches(0.22), font_size=9, bold=True, color=WHITE)

        add_multiline_text(slide,
                           [l for l in desc.split('\n')],
                           cx + Inches(0.1), cy + Inches(1.85),
                           Inches(2.8), Inches(2.8),
                           font_size=11, color=TEXT_MUTED, spacing=4)

    # Features used
    add_rect(slide, Inches(0.4), Inches(6.65), Inches(12.5), Inches(0.65), fill_color=BG_CARD)
    add_rect(slide, Inches(0.4), Inches(6.65), Inches(0.06), Inches(0.65), fill_color=BLUE)
    add_text(slide, "Features: ville · surface (m²) · nbre chambres · nbre SDB · type de bien · quartier   →   Prix (MAD)",
             Inches(0.6), Inches(6.72), Inches(12.0), Inches(0.45),
             font_size=12, color=TEXT_LIGHT)

    add_morph_transition(slide)
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — DASHBOARDS POWER BI
# ─────────────────────────────────────────────────────────────────────────────
def slide_dashboards():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_bg(slide, BG_DARK)
    section_header_bar(slide, "07 — Dashboards Analytiques — Power BI", TEAL)

    # Mock dashboard frame
    add_rect(slide, Inches(0.4), Inches(1.0), Inches(7.8), Inches(5.8), fill_color=BG_CARD)
    add_gradient_rect(slide, Inches(0.4), Inches(1.0), Inches(7.8), Inches(0.08),
                      TEAL, BLUE, angle=0)
    add_text(slide, "Power BI Embedded Dashboard", Inches(0.6), Inches(1.15),
             Inches(7.4), Inches(0.4), font_size=13, bold=True, color=TEAL)

    # Fake chart areas inside dashboard
    add_rect(slide, Inches(0.6), Inches(1.65), Inches(3.6), Inches(2.0), fill_color=BG_SURFACE)
    add_text(slide, "📈 Évolution des Prix\npar Ville", Inches(0.7), Inches(1.85),
             Inches(3.4), Inches(1.6), font_size=13, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(4.4), Inches(1.65), Inches(3.6), Inches(2.0), fill_color=BG_SURFACE)
    add_text(slide, "🥧 Répartition\npar Type de Bien", Inches(4.5), Inches(1.85),
             Inches(3.4), Inches(1.6), font_size=13, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(0.6), Inches(3.8), Inches(7.4), Inches(2.8), fill_color=BG_SURFACE)
    add_text(slide, "🗺️  Carte de Chaleur — Distribution des Prix par Région Marocaine",
             Inches(0.7), Inches(4.3), Inches(7.2), Inches(1.8),
             font_size=14, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # Right column: KPIs
    add_rect(slide, Inches(8.5), Inches(1.0), Inches(4.4), Inches(5.8), fill_color=BG_CARD)
    add_rect(slide, Inches(8.5), Inches(1.0), Inches(0.06), Inches(5.8), fill_color=ORANGE)
    add_text(slide, "📊  KPIs Clés", Inches(8.7), Inches(1.1),
             Inches(4.0), Inches(0.45), font_size=15, bold=True, color=ORANGE)

    kpis = [
        ("7 805", "Annonces totales", BLUE),
        ("331", "Villes couvertes", PURPLE),
        ("1 350 000 MAD", "Prix médian", TEAL),
        ("2 837 789 MAD", "Prix moyen", ORANGE),
        ("9", "Sources scrappées", YELLOW),
        ("20/h", "Estimations/utilisateur", RED),
    ]
    for i, (val, label, color) in enumerate(kpis):
        ky = Inches(1.65 + i * 0.82)
        add_rect(slide, Inches(8.7), ky, Inches(4.0), Inches(0.68), fill_color=BG_SURFACE)
        add_rect(slide, Inches(8.7), ky, Inches(0.04), Inches(0.68), fill_color=color)
        add_text(slide, val, Inches(8.85), ky + Inches(0.06),
                 Inches(2.5), Inches(0.35), font_size=16, bold=True, color=color)
        add_text(slide, label, Inches(8.85), ky + Inches(0.38),
                 Inches(3.7), Inches(0.22), font_size=10, color=TEXT_MUTED)

    add_morph_transition(slide)
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — CHATBOT N8N
# ─────────────────────────────────────────────────────────────────────────────
def slide_chatbot():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_bg(slide, BG_DARK)
    section_header_bar(slide, "08 — Chatbot Intelligent — N8N Integration", ORANGE)

    # Chat window mock
    add_rect(slide, Inches(0.4), Inches(1.0), Inches(5.5), Inches(6.2), fill_color=BG_CARD)
    add_gradient_rect(slide, Inches(0.4), Inches(1.0), Inches(5.5), Inches(0.55),
                      BLUE, PURPLE, angle=0)
    add_text(slide, "💬  Assistant ORCHIIMMO", Inches(0.6), Inches(1.08),
             Inches(5.0), Inches(0.38), font_size=14, bold=True, color=WHITE)
    add_text(slide, "● En ligne", Inches(0.6), Inches(1.42),
             Inches(3.0), Inches(0.18), font_size=9, color=TEAL)

    # Chat messages mock
    msgs = [
        ("bot", "Bonjour ! Comment puis-je vous aider avec l'immobilier au Maroc ?", BLUE),
        ("user", "Quel est le prix moyen d'un appartement à Casablanca ?", PURPLE),
        ("bot", "À Casablanca, le prix moyen est de ~2 500 000 MAD pour un appart de 85m². Je peux vous faire une estimation précise si vous le souhaitez !", BLUE),
        ("user", "Merci ! Pouvez-vous m'aider à estimer un bien ?", PURPLE),
        ("bot", "Bien sûr ! Rendez-vous sur la page Estimation IA pour obtenir une prédiction personnalisée avec intervalle de confiance.", BLUE),
    ]
    for i, (sender, msg, color) in enumerate(msgs):
        my = Inches(1.65 + i * 0.98)
        is_bot = sender == "bot"
        mx = Inches(0.5) if is_bot else Inches(1.5)
        mw = Inches(4.5)
        add_rect(slide, mx, my, mw, Inches(0.78), fill_color=BG_SURFACE)
        add_rect(slide, mx if is_bot else mx + mw - Inches(0.04),
                 my, Inches(0.04), Inches(0.78), fill_color=color)
        add_text(slide, msg, mx + Inches(0.12), my + Inches(0.1),
                 mw - Inches(0.2), Inches(0.62), font_size=9, color=TEXT_LIGHT)

    # Right: tech details
    add_rect(slide, Inches(6.3), Inches(1.0), Inches(6.6), Inches(6.2), fill_color=BG_CARD)
    add_rect(slide, Inches(6.3), Inches(1.0), Inches(0.06), Inches(6.2), fill_color=ORANGE)
    add_text(slide, "⚙️  Architecture Chatbot", Inches(6.5), Inches(1.1),
             Inches(6.0), Inches(0.45), font_size=15, bold=True, color=ORANGE)

    tech_blocks = [
        ("N8N Workflow Engine", "Orchestration des flux de conversation\nWebhooks + logique métier", ORANGE),
        ("Claude API (Anthropic)", "Modèle IA sous-jacent pour\nla compréhension du langage naturel", BLUE),
        ("Django Proxy /chatbot/", "Endpoint sécurisé qui relaie\nles requêtes vers N8N", PURPLE),
        ("Widget Frontend", "@n8n/chat CDN — s'intègre sur\ntoutes les pages du site", TEAL),
        ("Bilingue FR / AR", "Support Français et Arabe Darija\nRTL layout automatique", YELLOW),
    ]
    for i, (title, desc, color) in enumerate(tech_blocks):
        ty = Inches(1.65 + i * 1.0)
        add_rect(slide, Inches(6.5), ty, Inches(6.1), Inches(0.85), fill_color=BG_SURFACE)
        add_rect(slide, Inches(6.5), ty, Inches(0.04), Inches(0.85), fill_color=color)
        add_text(slide, title, Inches(6.65), ty + Inches(0.05),
                 Inches(5.8), Inches(0.32), font_size=12, bold=True, color=color)
        add_text(slide, desc, Inches(6.65), ty + Inches(0.38),
                 Inches(5.8), Inches(0.42), font_size=10, color=TEXT_MUTED)

    add_wipe_transition(slide, "r")
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — INTERFACE UI/UX
# ─────────────────────────────────────────────────────────────────────────────
def slide_uiux():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_bg(slide, BG_DARK)
    section_header_bar(slide, "Interface Utilisateur — Design System", BLUE)

    add_text(slide, "Design Glassmorphism · Dark / Light Mode · Bilingue FR / AR",
             Inches(0.65), Inches(0.75), Inches(10), Inches(0.35),
             font_size=13, color=TEXT_MUTED)

    # Color palette display
    colors_display = [
        (BLUE,   "#4895EF", "Bleu Accent"),
        (PURPLE, "#7C3AED", "Violet"),
        (TEAL,   "#06D6A0", "Teal"),
        (ORANGE, "#F4A261", "Orange"),
        (RED,    "#EF4444", "Rouge"),
        (YELLOW, "#F59E0B", "Jaune"),
    ]
    add_text(slide, "Palette de Couleurs", Inches(0.4), Inches(1.15),
             Inches(5), Inches(0.35), font_size=13, bold=True, color=TEXT_LIGHT)
    for i, (color, hex_val, name) in enumerate(colors_display):
        cx = Inches(0.4 + i * 1.05)
        cy = Inches(1.55)
        add_rect(slide, cx, cy, Inches(0.9), Inches(0.55), fill_color=color)
        add_text(slide, hex_val, cx, cy + Inches(0.6),
                 Inches(0.9), Inches(0.25), font_size=8, color=TEXT_MUTED,
                 align=PP_ALIGN.CENTER)
        add_text(slide, name, cx, cy + Inches(0.85),
                 Inches(0.9), Inches(0.25), font_size=8, color=TEXT_MUTED,
                 align=PP_ALIGN.CENTER)

    # UI features
    ui_features = [
        ("🌑 Dark Mode", "Thème sombre par défaut\n#080B14 fond, glassmorphism", BLUE),
        ("☀️ Light Mode", "Basculement dynamique\nLocalStorage persistance", YELLOW),
        ("🇲🇦 Bilingue", "Français ↔ Arabe Darija\nRTL automatique pour AR", TEAL),
        ("✨ Animations", "Tilt 3D, gradient shifts\nPulse-glow, shimmer effect", PURPLE),
        ("📱 Responsive", "Bootstrap 5.3 grid\nMobile / Tablet / Desktop", ORANGE),
        ("🔒 Sécurité", "Rate limiting 20/h\nCSRF, auth required", RED),
    ]
    for i, (title, desc, color) in enumerate(ui_features):
        row = i // 3
        col = i % 3
        cx = Inches(0.4 + col * 4.3)
        cy = Inches(2.7 + row * 2.1)
        cw = Inches(4.1)
        ch = Inches(1.9)

        add_rect(slide, cx, cy, cw, ch, fill_color=BG_CARD)
        add_rect(slide, cx, cy, Inches(0.06), ch, fill_color=color)
        add_text(slide, title, cx + Inches(0.15), cy + Inches(0.15),
                 cw - Inches(0.25), Inches(0.4), font_size=14, bold=True, color=color)
        add_text(slide, desc, cx + Inches(0.15), cy + Inches(0.6),
                 cw - Inches(0.25), Inches(1.2), font_size=11, color=TEXT_MUTED)

    add_morph_transition(slide)
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — DEPLOIEMENT
# ─────────────────────────────────────────────────────────────────────────────
def slide_deploiement():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_bg(slide, BG_DARK)
    section_header_bar(slide, "09 — Déploiement & Infrastructure Production", RED)

    # Architecture diagram (simplified)
    boxes = [
        (Inches(0.5), Inches(1.2), Inches(2.8), Inches(1.0),
         "👤 Utilisateur", "Navigateur web\nFR / AR", BLUE),
        (Inches(4.5), Inches(1.2), Inches(2.8), Inches(1.0),
         "☁️ Render.com", "Cloud Hosting\nFrankfurt (EU)", RED),
        (Inches(8.5), Inches(1.2), Inches(2.8), Inches(1.0),
         "⚙️ Django 5.2", "Gunicorn WSGI\nPython 3.12", PURPLE),
        (Inches(0.5), Inches(3.0), Inches(2.8), Inches(1.0),
         "🗄️ PostgreSQL", "Render DB\nFree Tier", TEAL),
        (Inches(4.5), Inches(3.0), Inches(2.8), Inches(1.0),
         "📦 WhiteNoise", "Static Files\nProduction", ORANGE),
        (Inches(8.5), Inches(3.0), Inches(2.8), Inches(1.0),
         "📧 Brevo", "Email Backend\nPassword Reset", YELLOW),
    ]
    for bx, by, bw, bh, title, desc, color in boxes:
        add_rect(slide, bx, by, bw, bh, fill_color=BG_CARD)
        add_gradient_rect(slide, bx, by, bw, Inches(0.05), color, PURPLE, angle=0)
        add_text(slide, title, bx + Inches(0.1), by + Inches(0.12),
                 bw - Inches(0.2), Inches(0.38), font_size=13, bold=True, color=color)
        add_text(slide, desc, bx + Inches(0.1), by + Inches(0.52),
                 bw - Inches(0.2), Inches(0.4), font_size=10, color=TEXT_MUTED)

    # Arrows between boxes
    for ax, ay, txt in [
        (Inches(3.4), Inches(1.65), "HTTPS →"),
        (Inches(7.4), Inches(1.65), "WSGI →"),
        (Inches(3.4), Inches(3.45), "SQL →"),
        (Inches(7.4), Inches(3.45), "SMTP →"),
    ]:
        add_text(slide, txt, ax, ay, Inches(1.0), Inches(0.3),
                 font_size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # Deployment config
    add_rect(slide, Inches(0.5), Inches(4.3), Inches(12.0), Inches(2.9), fill_color=BG_CARD)
    add_rect(slide, Inches(0.5), Inches(4.3), Inches(0.06), Inches(2.9), fill_color=RED)
    add_text(slide, "⚙️  Configuration render.yaml", Inches(0.7), Inches(4.4),
             Inches(11), Inches(0.4), font_size=14, bold=True, color=RED)

    config_lines = [
        {"text": "services:", "color": BLUE, "size": 11, "bold": True},
        {"text": "  type: web  |  env: python  |  region: frankfurt  |  plan: free", "color": TEXT_MUTED, "size": 10, "bold": False},
        {"text": "  buildCommand: pip install -r requirements.txt && python manage.py collectstatic", "color": TEXT_MUTED, "size": 10, "bold": False},
        {"text": "  startCommand: gunicorn config.wsgi:application", "color": TEAL, "size": 10, "bold": False},
        {"text": "databases:", "color": BLUE, "size": 11, "bold": True},
        {"text": "  type: pgsql  |  plan: free  |  region: frankfurt", "color": TEXT_MUTED, "size": 10, "bold": False},
        {"text": "URL Production:  orchiimmo.onrender.com", "color": ORANGE, "size": 12, "bold": True},
    ]
    add_multiline_text(slide, config_lines, Inches(0.7), Inches(4.88),
                       Inches(11.5), Inches(2.2), spacing=2)

    add_push_transition(slide, "l")
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — RESULTATS & BILAN
# ─────────────────────────────────────────────────────────────────────────────
def slide_resultats():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_bg(slide, BG_DARK)
    section_header_bar(slide, "10 — Résultats, Bilan & Perspectives", TEAL)

    # Big KPI row
    kpis = [
        ("7 805", "Annonces\ncollectées", BLUE),
        ("91.2%", "R² Score\nLightGBM", TEAL),
        ("331", "Villes\ncouvertes", PURPLE),
        ("9", "Sources\nscrappées", ORANGE),
    ]
    for i, (val, label, color) in enumerate(kpis):
        cx = Inches(0.4 + i * 3.15)
        add_rect(slide, cx, Inches(1.1), Inches(2.9), Inches(1.6), fill_color=BG_CARD)
        add_gradient_rect(slide, cx, Inches(1.1), Inches(2.9), Inches(0.06),
                          color, PURPLE, angle=0)
        add_text(slide, val, cx + Inches(0.1), Inches(1.22),
                 Inches(2.7), Inches(0.85), font_size=38, bold=True, color=color,
                 align=PP_ALIGN.CENTER)
        add_text(slide, label, cx + Inches(0.1), Inches(2.05),
                 Inches(2.7), Inches(0.55), font_size=11, color=TEXT_MUTED,
                 align=PP_ALIGN.CENTER)

    # Accomplishments
    add_rect(slide, Inches(0.4), Inches(2.95), Inches(6.0), Inches(4.3), fill_color=BG_CARD)
    add_rect(slide, Inches(0.4), Inches(2.95), Inches(0.06), Inches(4.3), fill_color=TEAL)
    add_text(slide, "✅  Réalisations du Stage", Inches(0.6), Inches(3.05),
             Inches(5.6), Inches(0.4), font_size=15, bold=True, color=TEAL)

    achievements = [
        "✅  Plateforme web complète déployée en production",
        "✅  Moteur IA avec 4 modèles ML comparés",
        "✅  Web scraping de 9 portails immobiliers",
        "✅  7 805 annonces collectées sur 331 villes",
        "✅  Dashboards Power BI intégrés",
        "✅  Chatbot N8N en FR et AR (Darija)",
        "✅  Interface glassmorphism dark/light bilingue",
        "✅  Déploiement cloud Render.com (Frankfurt)",
    ]
    add_multiline_text(slide,
                       [{"text": a, "color": TEXT_LIGHT, "size": 12, "bold": False}
                        for a in achievements],
                       Inches(0.6), Inches(3.55), Inches(5.7), Inches(3.5), spacing=3)

    # Perspectives
    add_rect(slide, Inches(6.7), Inches(2.95), Inches(6.2), Inches(4.3), fill_color=BG_CARD)
    add_rect(slide, Inches(6.7), Inches(2.95), Inches(0.06), Inches(4.3), fill_color=ORANGE)
    add_text(slide, "🚀  Perspectives d'Amélioration", Inches(6.9), Inches(3.05),
             Inches(5.8), Inches(0.4), font_size=15, bold=True, color=ORANGE)

    perspectives = [
        "🔮  Application mobile React Native",
        "🔮  Modèles deep learning (Neural Networks)",
        "🔮  Alertes email/SMS sur nouveaux biens",
        "🔮  Comparateur de biens côte à côte",
        "🔮  API publique REST pour partenaires",
        "🔮  Intégration de photos pour estimation",
        "🔮  Score de quartier (écoles, commerces...)",
        "🔮  Analyse prédictive des tendances",
    ]
    add_multiline_text(slide,
                       [{"text": p, "color": TEXT_MUTED, "size": 12, "bold": False}
                        for p in perspectives],
                       Inches(6.9), Inches(3.55), Inches(5.8), Inches(3.5), spacing=3)

    add_morph_transition(slide)
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15 — CONCLUSION
# ─────────────────────────────────────────────────────────────────────────────
def slide_conclusion():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_bg(slide, BG_DARK)
    add_gradient_rect(slide, 0, 0, W, Inches(0.08), BLUE, PURPLE, angle=0)

    # Big title
    add_text(slide, "Conclusion", Inches(0.5), Inches(0.3),
             Inches(12), Inches(0.7), font_size=36, bold=True, color=TEXT_LIGHT)
    add_gradient_rect(slide, Inches(0.5), Inches(0.95), Inches(3.5), Inches(0.05),
                      BLUE, PURPLE, angle=0)

    # Summary paragraph
    summary = (
        "Ce projet de stage de fin d'études a permis de concevoir et déployer "
        "ORCHIIMMO, une plateforme immobilière complète dédiée au marché marocain. "
        "En combinant le web scraping de 9 sources, un moteur IA basé sur LightGBM "
        "(R² = 91.2%), des dashboards Power BI et un chatbot intelligent N8N, "
        "la plateforme répond à un réel besoin de transparence et d'accessibilité "
        "des données immobilières au Maroc."
    )
    add_text(slide, summary, Inches(0.5), Inches(1.1), Inches(12.3), Inches(1.5),
             font_size=14, color=TEXT_MUTED)

    # Three pillars
    pillars = [
        ("🔬", "Innovation\nTechnique",
         "4 modèles ML comparés\nLightGBM sélectionné\nR² Score de 91.2%", BLUE),
        ("🌐", "Impact\nMétier",
         "7 805 annonces collectées\n331 villes marocaines\nPrix transparents et accessibles", TEAL),
        ("🚀", "Déploiement\nProfessionnel",
         "Cloud Render.com (Frankfurt)\nCI/CD avec render.yaml\norchiimmo.onrender.com", PURPLE),
    ]
    for i, (icon, title, desc, color) in enumerate(pillars):
        cx = Inches(0.5 + i * 4.1)
        cy = Inches(2.85)
        add_rect(slide, cx, cy, Inches(3.8), Inches(3.5), fill_color=BG_CARD)
        add_gradient_rect(slide, cx, cy, Inches(3.8), Inches(0.06), color, PURPLE, angle=0)

        ic = slide.shapes.add_shape(9, cx + Inches(1.15), cy + Inches(0.2),
                                     Inches(1.5), Inches(1.5))
        ic.fill.solid()
        ic.fill.fore_color.rgb = BG_SURFACE
        ic.line.color.rgb = color
        ic.line.width = Pt(1.5)

        add_text(slide, icon, cx + Inches(1.15), cy + Inches(0.28),
                 Inches(1.5), Inches(1.2), font_size=36, align=PP_ALIGN.CENTER)
        add_text(slide, title, cx + Inches(0.1), cy + Inches(1.8),
                 Inches(3.6), Inches(0.55), font_size=14, bold=True,
                 color=color, align=PP_ALIGN.CENTER)
        add_text(slide, desc, cx + Inches(0.1), cy + Inches(2.45),
                 Inches(3.6), Inches(1.0), font_size=11, color=TEXT_MUTED,
                 align=PP_ALIGN.CENTER)

    add_text(slide, "ORCHIIMMO © 2026 — Plateforme Immobilière Maroc | PFE 2025-2026",
             Inches(0.5), Inches(6.95), Inches(12), Inches(0.35),
             font_size=10, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    add_fade_transition(slide)
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 16 — MERCI / QUESTIONS
# ─────────────────────────────────────────────────────────────────────────────
def slide_merci():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    set_bg(slide, BG_DARK)

    # Full gradient background
    add_gradient_rect(slide, 0, 0, W, H, BG_DARK, BG_CARD, angle=135)
    add_gradient_rect(slide, 0, 0, W, Inches(0.12), BLUE, PURPLE, angle=0)
    add_gradient_rect(slide, 0, H - Inches(0.12), W, Inches(0.12), PURPLE, TEAL, angle=0)

    # Decorative circles
    for i, (color, size, x, y) in enumerate([
        (BLUE, 4.0, 11.5, -1.5),
        (PURPLE, 2.5, 0.5, 5.5),
        (TEAL, 1.5, 10.5, 5.5),
    ]):
        circ = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(size), Inches(size))
        circ.fill.solid()
        circ.fill.fore_color.rgb = color
        circ.line.fill.background()
        sp = circ._element
        # Add transparency via XML
        spPr = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')
        if spPr is not None:
            solidFill = spPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
            if solidFill is not None:
                srgb = solidFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                if srgb is not None:
                    alpha_elem = etree.SubElement(
                        srgb,
                        '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha'
                    )
                    alpha_elem.set('val', '8000')

    # ORCHI IMMO big
    add_text(slide, "ORCHI", Inches(2.0), Inches(1.2), Inches(5), Inches(1.5),
             font_size=80, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, "IMMO", Inches(6.5), Inches(1.2), Inches(5.5), Inches(1.5),
             font_size=80, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

    add_gradient_rect(slide, Inches(2.0), Inches(2.6), Inches(9.5), Inches(0.06),
                      BLUE, PURPLE, angle=0)

    add_text(slide, "Merci pour votre attention !",
             Inches(1.5), Inches(2.75), Inches(10.5), Inches(0.7),
             font_size=28, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    add_text(slide, "Place aux Questions 🎤",
             Inches(1.5), Inches(3.5), Inches(10.5), Inches(0.55),
             font_size=22, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # Bottom info
    add_rect(slide, Inches(2.0), Inches(4.5), Inches(9.5), Inches(1.8), fill_color=BG_CARD)
    add_gradient_rect(slide, Inches(2.0), Inches(4.5), Inches(9.5), Inches(0.04),
                      BLUE, PURPLE, angle=0)

    contact_items = [
        ("📧", "hadifatimaezzahra5@gmail.com", BLUE),
        ("🌐", "orchiimmo.onrender.com", TEAL),
        ("🇲🇦", "PFE 2025-2026 | Maroc", ORANGE),
    ]
    for i, (icon, text, color) in enumerate(contact_items):
        add_text(slide, f"{icon}  {text}",
                 Inches(2.5 + i * 3.2), Inches(4.75),
                 Inches(3.0), Inches(0.45),
                 font_size=13, color=color, align=PP_ALIGN.CENTER)

    add_text(slide, "ORCHIIMMO © 2026 — Plateforme Immobilière Maroc | PFE 2025-2026",
             Inches(0.5), Inches(6.95), Inches(12), Inches(0.35),
             font_size=10, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    add_morph_transition(slide)
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# BUILD ALL SLIDES
# ─────────────────────────────────────────────────────────────────────────────
print("🏗️  Construction de la présentation ORCHIIMMO...")
slide_cover()
print("  ✅ Slide 1 — Couverture")
slide_plan()
print("  ✅ Slide 2 — Plan")
slide_contexte()
print("  ✅ Slide 3 — Contexte & Problématique")
slide_projet()
print("  ✅ Slide 4 — Présentation Projet")
slide_architecture()
print("  ✅ Slide 5 — Architecture Technique")
slide_scraping()
print("  ✅ Slide 6 — Scraping & Données")
slide_catalogue()
print("  ✅ Slide 7 — Catalogue & Carte")
slide_estimation()
print("  ✅ Slide 8 — Estimation IA")
slide_ml()
print("  ✅ Slide 9 — Modèles ML")
slide_dashboards()
print("  ✅ Slide 10 — Dashboards Power BI")
slide_chatbot()
print("  ✅ Slide 11 — Chatbot N8N")
slide_uiux()
print("  ✅ Slide 12 — Interface UI/UX")
slide_deploiement()
print("  ✅ Slide 13 — Déploiement")
slide_resultats()
print("  ✅ Slide 14 — Résultats & Perspectives")
slide_conclusion()
print("  ✅ Slide 15 — Conclusion")
slide_merci()
print("  ✅ Slide 16 — Merci / Questions")

# Save
output_path = "/home/user/ORCHIIMMO/ORCHIIMMO_Soutenance_SFE_2026.pptx"
prs.save(output_path)
print(f"\n🎉  Présentation sauvegardée: {output_path}")
print(f"📊  Total slides: {len(prs.slides)}")
